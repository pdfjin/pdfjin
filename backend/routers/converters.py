from fastapi import APIRouter, File, UploadFile, HTTPException, Form, Response
from typing import List
import io
import os
import subprocess
import tempfile
import zipfile
import img2pdf
from pypdf import PdfReader, PdfWriter
from pdf2docx import Converter
from pdf2image import convert_from_bytes
from pptx import Presentation
import ocrmypdf
import pikepdf

router = APIRouter()

# ─── TOOL: PDF TO EXCEL ───────────────────────────────────────
@router.post("/pdf-to-excel")
async def pdf_to_excel(files: List[UploadFile] = File(...)):
    import pdfplumber
    import pandas as pd
    file = files[0]
    try:
        from database import load_db, save_db
        try:
            db = load_db()
            db["stats"]["conversions_today"] = db.get("stats", {}).get("conversions_today", 0) + 1
            save_db(db)
        except: pass
        
        content = await file.read()
        all_tables = []
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for i, page in enumerate(pdf.pages):
                tables = page.extract_tables()
                for table_idx, raw_table in enumerate(tables):
                    if not raw_table: continue
                    df = pd.DataFrame(raw_table)
                    all_tables.append((f"Page{i+1}_Table{table_idx+1}", df))
        
        if not all_tables:
            raise Exception("No tables detected in the PDF.")

        output = io.BytesIO()
        with pd.ExcelWriter(output, engine='openpyxl') as writer:
            for sheet_name, df in all_tables:
                # Clean sheet name to be excel compatible
                clean_name = "".join([c for c in sheet_name if c.isalnum() or c in " _-"])[:31]
                df.to_excel(writer, sheet_name=clean_name, index=False, header=False)
        output.seek(0)
        
        base_name = file.filename.rsplit('.', 1)[0] if file.filename else "converted"
        return Response(
            content=output.getvalue(), 
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", 
            headers={"Content-Disposition": f"attachment; filename={base_name}_pdfjin.xlsx"}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Excel conversion failed: {str(e)}")

# ─── TOOL: WORD TO PDF ────────────────────────────────────────
@router.post("/word-to-pdf")
async def word_to_pdf(files: List[UploadFile] = File(...)):
    file = files[0]
    try:
        from database import load_db, save_db
        try:
            db = load_db()
            if "stats" not in db: db["stats"] = {}
            db["stats"]["conversions_today"] = db["stats"].get("conversions_today", 0) + 1
            save_db(db)
        except: pass

        with tempfile.TemporaryDirectory() as temp_dir:
            orig_name = file.filename or "document.docx"
            input_path = os.path.join(temp_dir, orig_name)
            with open(input_path, "wb") as f:
                f.write(await file.read())
            
            # Use libreoffice to convert
            subprocess.run(
                ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', temp_dir, input_path],
                check=True, timeout=60
            )
            
            # Find generated PDF
            pdf_path = None
            for f in os.listdir(temp_dir):
                if f.lower().endswith(".pdf"):
                    pdf_path = os.path.join(temp_dir, f)
                    pdf_filename = f
                    break
            
            if not pdf_path or not os.path.exists(pdf_path):
                raise Exception("PDF generation failed.")

            with open(pdf_path, "rb") as f:
                pdf_data = f.read()

            return Response(
                content=pdf_data, 
                media_type="application/pdf", 
                headers={
                    "Content-Disposition": f"attachment; filename={pdf_filename}",
                    "Access-Control-Expose-Headers": "Content-Disposition"
                }
            )
    except Exception as e:
        print(f"WORD2PDF ERROR: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Word to PDF failed: {str(e)}")

# ─── TOOL: PDF TO WORD ────────────────────────────────────────
@router.post("/pdf-to-word")
async def pdf_to_word(files: List[UploadFile] = File(...)):
    file = files[0]
    try:
        # Update Stats (Nuclear Update for "Server is updating files")
        from database import load_db, save_db
        try:
            db = load_db()
            if "stats" not in db: db["stats"] = {}
            db["stats"]["conversions_today"] = db["stats"].get("conversions_today", 0) + 1
            save_db(db)
        except Exception as se:
            print(f"Stats update failed: {se}")

        with tempfile.TemporaryDirectory() as temp_dir:
            pdf_in = os.path.join(temp_dir, "input.pdf")
            pdf_clean = os.path.join(temp_dir, "clean.pdf")
            doc_out = os.path.join(temp_dir, "output.docx")
            
            # 1. Save upload to disk
            content = await file.read()
            with open(pdf_in, "wb") as f:
                f.write(content)
            
            # 2. Sanitize with pikepdf to fix common issues
            try:
                with pikepdf.open(pdf_in) as p:
                    p.save(pdf_clean)
                convert_src = pdf_clean
            except:
                convert_src = pdf_in # Fallback to original if pikepdf fails
                
            # 3. Perform Conversion
            cv = Converter(convert_src)
            cv.convert(doc_out, start=0, end=None)
            cv.close()
            
            # 4. Read result
            if not os.path.exists(doc_out):
                raise Exception("DOCX generation failed.")
                
            with open(doc_out, "rb") as f:
                docx_data = f.read()
            
            # 5. Clean filename
            orig_name = file.filename or "document.pdf"
            clean_name = orig_name.rsplit('.', 1)[0]
            final_name = f"{clean_name}_pdfjin.docx"
            
            return Response(
                content=docx_data,
                media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                headers={
                    'Content-Disposition': f'attachment; filename="{final_name}"',
                    'Access-Control-Expose-Headers': 'Content-Disposition'
                }
            )
    except Exception as e:
        print(f"PDF2WORD ERROR: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Conversion Error: {str(e)}")

# ─── TOOL: JPG TO PDF ─────────────────────────────────────────
@router.post("/jpg-to-pdf")
async def jpg_to_pdf(files: List[UploadFile] = File(...)):
    try:
        from database import load_db, save_db
        try:
            db = load_db()
            if "stats" not in db: db["stats"] = {}
            db["stats"]["conversions_today"] = db["stats"].get("conversions_today", 0) + 1
            save_db(db)
        except: pass

        from PIL import Image
        processed_images = []
        
        for f in files:
            content = await f.read()
            try:
                img = Image.open(io.BytesIO(content))
                # Normalize to RGB (img2pdf works best with RGB/CMYK JPEGs)
                # This also handles PNG with Alpha and WEBP
                if img.mode in ("RGBA", "P", "LA") or (img.mode != "RGB" and img.mode != "CMYK"):
                    img = img.convert("RGB")
                    
                # We convert to JPEG bytes to ensure img2pdf can handle it losslessly where possible,
                # or as a standard format.
                buf = io.BytesIO()
                img.save(buf, format="JPEG", quality=95)
                processed_images.append(buf.getvalue())
            except Exception as img_err:
                print(f"Error processing image {f.filename}: {img_err}")
                # If Pillow fails, try passing raw content as fallback
                processed_images.append(content)
        
        if not processed_images:
            raise Exception("No valid images provided.")

        pdf_bytes = img2pdf.convert(processed_images)
        return Response(
            content=pdf_bytes, 
            media_type="application/pdf", 
            headers={
                "Content-Disposition": "attachment; filename=images_to_pdfjin.pdf",
                "Access-Control-Expose-Headers": "Content-Disposition"
            }
        )
    except Exception as e:
        print(f"JPG2PDF ERROR: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Image conversion failed: {str(e)}")

# ─── TOOL: PDF TO JPG ─────────────────────────────────────────
@router.post("/pdf-to-jpg")
async def pdf_to_jpg(files: List[UploadFile] = File(...)):
    file = files[0]
    try:
        from database import load_db, save_db
        try:
            db = load_db()
            db["stats"]["conversions_today"] = db.get("stats", {}).get("conversions_today", 0) + 1
            save_db(db)
        except: pass

        content = await file.read()
        images = convert_from_bytes(content)
        
        if len(images) == 1:
            buf = io.BytesIO()
            images[0].save(buf, format='JPEG', quality=95)
            buf.seek(0)
            return Response(content=buf.getvalue(), media_type="image/jpeg", headers={"Content-Disposition": f"attachment; filename={file.filename.rsplit('.', 1)[0]}.jpg"})
        
        zip_buf = io.BytesIO()
        with zipfile.ZipFile(zip_buf, "w") as zf:
            for i, img in enumerate(images):
                buf = io.BytesIO()
                img.save(buf, format='JPEG', quality=95)
                zf.writestr(f"page_{i+1}.jpg", buf.getvalue())
        
        zip_buf.seek(0)
        return Response(content=zip_buf.getvalue(), media_type="application/zip", headers={"Content-Disposition": f"attachment; filename={file.filename.rsplit('.', 1)[0]}_images.zip"})
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF to JPG failed: {str(e)}")

# ─── TOOL: OCR PDF ────────────────────────────────────────────
@router.post("/ocr-pdf")
async def ocr_pdf(files: List[UploadFile] = File(...)):
    file = files[0]
    print(f"OCR: Processing {file.filename}...")
    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            in_p = os.path.join(temp_dir, "in.pdf")
            out_p = os.path.join(temp_dir, "out.pdf")
            
            content = await file.read()
            with open(in_p, "wb") as f:
                f.write(content)
            
            print(f"OCR: Starting ocrmypdf on {in_p}...")
            # Use skip_text=True to avoid re-OCRing if text exists, 
            # rotate_pages=True and deskew=True for better quality.
            try:
                ocrmypdf.ocr(in_p, out_p, skip_text=True, rotate_pages=True, deskew=True, optimize=1)
                print(f"OCR: Finished successfully.")
            except Exception as ocr_err:
                print(f"OCR: ocrmypdf failed: {str(ocr_err)}")
                raise Exception(f"OCR engine error: {str(ocr_err)}")
            
            if not os.path.exists(out_p):
                print(f"OCR: Output file {out_p} not found.")
                raise Exception("Output file was not generated.")
                
            with open(out_p, "rb") as f:
                data = f.read()
            
            return Response(
                content=data, 
                media_type="application/pdf", 
                headers={
                    "Content-Disposition": f"attachment; filename=ocr_{file.filename}",
                    "Access-Control-Expose-Headers": "Content-Disposition"
                }
            )
    except Exception as e:
        print(f"OCR ERROR: {str(e)}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"OCR failed: {str(e)}")

# ─── TOOL: POWERPOINT TO PDF ──────────────────────────────────
@router.post("/powerpoint-to-pdf")
async def ppt_to_pdf(files: List[UploadFile] = File(...)):
    file = files[0]
    try:
        from database import load_db, save_db
        try:
            db = load_db()
            if "stats" not in db: db["stats"] = {}
            db["stats"]["conversions_today"] = db["stats"].get("conversions_today", 0) + 1
            save_db(db)
        except: pass

        with tempfile.TemporaryDirectory() as temp_dir:
            orig_name = file.filename or "presentation.pptx"
            input_path = os.path.join(temp_dir, orig_name)
            with open(input_path, "wb") as f: f.write(await file.read())
            
            subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', temp_dir, input_path], check=True, timeout=60)
            
            # Find generated PDF
            pdf_path = None
            pdf_filename = orig_name.rsplit('.', 1)[0] + ".pdf"
            for f in os.listdir(temp_dir):
                if f.lower().endswith(".pdf"):
                    pdf_path = os.path.join(temp_dir, f)
                    pdf_filename = f
                    break
            
            if not pdf_path or not os.path.exists(pdf_path):
                raise Exception("PDF generation failed.")

            with open(pdf_path, "rb") as f:
                pdf_data = f.read()
            
            return Response(
                content=pdf_data, 
                media_type="application/pdf", 
                headers={
                    "Content-Disposition": f"attachment; filename={pdf_filename}",
                    "Access-Control-Expose-Headers": "Content-Disposition"
                }
            )
    except Exception as e:
        print(f"PPT2PDF ERROR: {str(e)}")
        raise HTTPException(status_code=500, detail=f"PPT conversion failed: {str(e)}")

# ─── TOOL: PDF TO POWERPOINT ──────────────────────────────────
@router.post("/pdf-to-powerpoint")
async def pdf_to_ppt(files: List[UploadFile] = File(...)):
    from pdf2image import convert_from_bytes
    file = files[0]
    try:
        images = convert_from_bytes(await file.read())
        prs = Presentation()
        # Remove default slides
        while len(prs.slides) > 0:
            # Slides don't have a clear way to delete all at once in python-pptx easily without re-init
            break 

        for img in images:
            with tempfile.NamedTemporaryFile(suffix=".jpg", delete=False) as tmp:
                img.save(tmp.name, format="JPEG")
                slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank slide
                # Resize slide to match image aspect ratio if possible, but default is fine for now
                slide.shapes.add_picture(tmp.name, 0, 0, width=prs.slide_width, height=prs.slide_height)
                os.unlink(tmp.name)
        
        output = io.BytesIO()
        prs.save(output); output.seek(0)
        return Response(content=output.getvalue(), media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers={"Content-Disposition": f"attachment; filename={file.filename.rsplit('.', 1)[0]}.pptx"})
    except Exception as e: raise HTTPException(status_code=500, detail=f"PDF to PPT failed: {str(e)}")

# ─── TOOL: EXCEL TO PDF ───────────────────────────────────────
@router.post("/excel-to-pdf")
async def excel_to_pdf(files: List[UploadFile] = File(...)):
    file = files[0]
    try:
        from database import load_db, save_db
        try:
            db = load_db()
            if "stats" not in db: db["stats"] = {}
            db["stats"]["conversions_today"] = db["stats"].get("conversions_today", 0) + 1
            save_db(db)
        except: pass

        with tempfile.TemporaryDirectory() as temp_dir:
            orig_name = file.filename or "spreadsheet.xlsx"
            input_path = os.path.join(temp_dir, orig_name)
            with open(input_path, "wb") as f: f.write(await file.read())
            
            subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', temp_dir, input_path], check=True, timeout=60)
            
            # Find generated PDF
            pdf_path = None
            pdf_filename = orig_name.rsplit('.', 1)[0] + ".pdf"
            for f in os.listdir(temp_dir):
                if f.lower().endswith(".pdf"):
                    pdf_path = os.path.join(temp_dir, f)
                    pdf_filename = f
                    break
            
            if not pdf_path or not os.path.exists(pdf_path):
                raise Exception("PDF generation failed.")

            with open(pdf_path, "rb") as f:
                pdf_data = f.read()
            
            return Response(
                content=pdf_data, 
                media_type="application/pdf", 
                headers={
                    "Content-Disposition": f"attachment; filename={pdf_filename}",
                    "Access-Control-Expose-Headers": "Content-Disposition"
                }
            )
    except Exception as e:
        print(f"EXCEL2PDF ERROR: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Excel conversion failed: {str(e)}")
# ─── TOOL: HTML TO PDF ────────────────────────────────────────
@router.post("/html-to-pdf")
async def html_to_pdf(files: List[UploadFile] = File(...)):
    file = files[0]
    try:
        from database import load_db, save_db
        try:
            db = load_db()
            if "stats" not in db: db["stats"] = {}
            db["stats"]["conversions_today"] = db["stats"].get("conversions_today", 0) + 1
            save_db(db)
        except: pass

        with tempfile.TemporaryDirectory() as temp_dir:
            orig_name = file.filename or "document.html"
            input_path = os.path.join(temp_dir, orig_name)
            with open(input_path, "wb") as f: f.write(await file.read())
            
            # Using libreoffice for html-to-pdf conversion
            subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', temp_dir, input_path], check=True, timeout=60)
            
            # Find generated PDF
            pdf_path = None
            pdf_filename = orig_name.rsplit('.', 1)[0] + ".pdf"
            for f in os.listdir(temp_dir):
                if f.lower().endswith(".pdf"):
                    pdf_path = os.path.join(temp_dir, f)
                    pdf_filename = f
                    break
            
            if not pdf_path or not os.path.exists(pdf_path):
                raise Exception("PDF generation failed.")

            with open(pdf_path, "rb") as f:
                pdf_data = f.read()
            
            return Response(
                content=pdf_data, 
                media_type="application/pdf", 
                headers={
                    "Content-Disposition": f"attachment; filename={pdf_filename}",
                    "Access-Control-Expose-Headers": "Content-Disposition"
                }
            )
    except Exception as e:
        print(f"HTML2PDF ERROR: {str(e)}")
        raise HTTPException(status_code=500, detail=f"HTML conversion failed: {str(e)}")
